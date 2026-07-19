from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "KawanDigi_Project_Charter.pptx"
ANDROID = ROOT / "public/images/android"
WEB = ROOT / "public/images/web-lapak-employee"

NAVY = RGBColor(43, 12, 90)
NAVY2 = RGBColor(58, 19, 111)
PINK = RGBColor(232, 62, 140)
PINK2 = RGBColor(255, 111, 181)
PURPLE = RGBColor(124, 58, 237)
BG = RGBColor(247, 244, 251)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(36, 23, 53)
MUTED = RGBColor(111, 102, 128)
LINE = RGBColor(232, 225, 242)
GREEN = RGBColor(23, 138, 87)
ORANGE = RGBColor(183, 121, 31)
RED = RGBColor(184, 50, 90)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, radius=True, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def textbox(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
            align=PP_ALIGN.LEFT, font="Aptos", valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    p.font.name = font; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    return box


def rich_text(slide, runs, x, y, w, h, size=18, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for text, bold, color in runs:
        r = p.add_run(); r.text = text; r.font.name = "Aptos"; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=15, color=TEXT, spacing=7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.level = 0; p.font.name = "Aptos"; p.font.size = Pt(size)
        p.font.color.rgb = color; p.space_after = Pt(spacing)
        p.text = "•  " + item
    return box


def add_image_contain(slide, path, x, y, w, h, bg=WHITE):
    rect(slide, x, y, w, h, bg, True, LINE)
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min((w - .12) / iw, (h - .12) / ih)
    pw, ph = iw * scale, ih * scale
    slide.shapes.add_picture(str(path), Inches(x + (w-pw)/2), Inches(y + (h-ph)/2),
                             width=Inches(pw), height=Inches(ph))


def add_image_cover(slide, path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    frame_ratio, img_ratio = w/h, iw/ih
    if img_ratio > frame_ratio:
        pic_h = h; pic_w = h * img_ratio
    else:
        pic_w = w; pic_h = w / img_ratio
    pic = slide.shapes.add_picture(str(path), Inches(x-(pic_w-w)/2), Inches(y-(pic_h-h)/2),
                                   width=Inches(pic_w), height=Inches(pic_h))
    # crop instead of relying on overflow
    if pic_w > w:
        crop = (pic_w-w)/pic_w/2; pic.crop_left = crop; pic.crop_right = crop
        pic.left = Inches(x); pic.width = Inches(w)
    if pic_h > h:
        crop = (pic_h-h)/pic_h/2; pic.crop_top = crop; pic.crop_bottom = crop
        pic.top = Inches(y); pic.height = Inches(h)
    return pic


def base_slide(title, kicker=None, dark=False):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.333, 7.5, NAVY if dark else BG, False)
    if not dark:
        rect(slide, 0, 0, .12, 7.5, PINK, False)
    if kicker:
        textbox(slide, kicker.upper(), .55, .32, 5.8, .28, 10, PINK if not dark else PINK2, True)
    textbox(slide, title, .55, .62, 12.1, .62, 27, WHITE if dark else NAVY, True)
    return slide


def footer(slide, n, label="PROJECT CHARTER · KAWANDIGI"):
    textbox(slide, label, .55, 7.15, 6, .18, 8, MUTED, True)
    textbox(slide, f"{n:02d}", 12.15, 7.12, .6, .2, 9, MUTED, True, PP_ALIGN.RIGHT)


def card(slide, title, body, x, y, w, h, accent=PINK, title_size=16, body_size=12):
    rect(slide, x, y, w, h, WHITE, True, LINE)
    rect(slide, x, y, .07, h, accent, False, accent)
    textbox(slide, title, x+.25, y+.2, w-.45, .35, title_size, NAVY, True)
    if isinstance(body, list):
        add_bullets(slide, body, x+.25, y+.7, w-.5, h-.85, body_size, MUTED, 4)
    else:
        textbox(slide, body, x+.25, y+.72, w-.5, h-.9, body_size, MUTED)


# 1 — Cover
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, NAVY, False)
rect(s, 8.9, -1.1, 5.5, 5.5, NAVY2, True, NAVY2)
rect(s, 10.1, 4.8, 4.0, 3.0, PINK, True, PINK)
rect(s, .65, .55, .58, .58, PINK, True, PINK)
textbox(s, "KD", .65, .67, .58, .25, 14, WHITE, True, PP_ALIGN.CENTER)
textbox(s, "KAWANDIGI · KOPERASI KARYAWAN DIGITAL", 1.4, .68, 6.2, .25, 11, WHITE, True)
textbox(s, "PROJECT CHARTER", .68, 1.65, 4.3, .3, 12, PINK2, True)
textbox(s, "Inisiasi Koperasi KawanDigi\n- Brawijaya Hospital", .68, 2.05, 8.2, 1.7, 32, WHITE, True)
textbox(s, "Fondasi layanan keuangan anggota dan unit usaha digital\nuntuk Brawijaya Hospital Antasari", .7, 4.05, 7.1, .8, 17, RGBColor(225,216,239))
for i, label in enumerate(["Pilot: Brawijaya Antasari", "Tahap: Inisiasi & MVP", "Platform: Web Mobile"]):
    rect(s, .7 + i*2.55, 5.35, 2.35, .48, RGBColor(75, 40, 125), True, RGBColor(105,70,150))
    textbox(s, label, .82 + i*2.55, 5.49, 2.1, .16, 9, WHITE, True, PP_ALIGN.CENTER)
textbox(s, "Prepared for project alignment & approval", .7, 6.82, 5.4, .2, 9, RGBColor(190,180,210))

# 2 — Executive summary
s = base_slide("Executive Summary", "01 · Overview")
textbox(s, "WHY", .65, 1.48, 1.0, .28, 11, PINK, True)
textbox(s, "Brawijaya Hospital belum memiliki koperasi dan platform digital terpadu untuk mengelola kesejahteraan anggota.", .65, 1.82, 5.7, 1.1, 21, NAVY, True)
card(s, "Solusi", "Koperasi Digital KawanDigi berbasis web mobile dengan Digital Saving Account, pembayaran, laporan keuangan, dan Lapak Employee.", 6.75, 1.45, 5.8, 1.55, PURPLE)
card(s, "Pilot", ["Brawijaya Hospital Antasari", "Validasi melalui UAT", "Implementasi terbatas & evaluasi"], .65, 3.35, 3.8, 2.25, PINK)
card(s, "Nilai Utama", ["Transparansi transaksi", "Efisiensi layanan anggota", "Fondasi unit usaha digital"], 4.75, 3.35, 3.8, 2.25, PURPLE)
card(s, "Outcome", ["MVP siap pilot", "Data dapat ditelusuri", "Adopsi ≥ 90%"], 8.85, 3.35, 3.7, 2.25, NAVY2)
footer(s, 2)

# 3 — Business case
s = base_slide("Business Case", "02 · Why & What")
card(s, "Pernyataan Masalah", "Saat ini Brawijaya Hospital memiliki 0 koperasi karyawan, 0 sistem terpusat untuk pengelolaan anggota dan transaksi, serta 0 kanal marketplace internal. Dampaknya, layanan kesejahteraan dan aktivitas keuangan anggota belum terkelola secara transparan atau terukur, sementara potensi pendapatan unit usaha internal belum dimanfaatkan.", .6, 1.45, 5.95, 3.95, PINK, 17, 14)
card(s, "Deskripsi Proyek (WHAT)", "Membangun ekosistem web mobile yang menyatukan keanggotaan, Digital Saving Account, pembayaran, laporan keuangan, dan Lapak Employee. Delivery dilakukan lintas unit oleh HR, Finance & Accounting, IT, Legal/Compliance, calon pengurus koperasi, perwakilan karyawan, dan tim KawanDigi.", 6.8, 1.45, 5.9, 1.7, PURPLE, 17, 12)
card(s, "Objektif SMART", "Dalam 4 bulan sejak Project Charter disetujui, menyelesaikan 100% fitur MVP sesuai scope, meluluskan 100% skenario UAT kritikal, dan memastikan minimal 90% pengguna pilot Brawijaya Hospital Antasari berhasil login serta mengakses dashboard.", 6.8, 3.35, 5.9, 2.05, NAVY2, 17, 13)
footer(s, 3)

# 4 — Deliverables & success
s = base_slide("Deliverables & Success Metrics", "03 · Measurable Outcome")
deliverables = ["Registrasi & login anggota", "Dashboard anggota", "Digital Saving Account", "Pembayaran anggota", "Laporan keuangan & SHU", "Lapak Employee"]
for i, item in enumerate(deliverables):
    x = .65 + (i%3)*2.05; y = 1.48 + (i//3)*1.0
    rect(s, x, y, 1.82, .78, WHITE, True, LINE)
    textbox(s, str(i+1), x+.12, y+.17, .34, .28, 13, PINK, True, PP_ALIGN.CENTER)
    textbox(s, item, x+.52, y+.14, 1.18, .44, 11, NAVY, True)
metrics = [("100%", "Fitur MVP sesuai scope"), ("100%", "Skenario UAT kritikal"), ("100%", "Transaksi pilot digital"), ("≥ 90%", "Login & akses dashboard")]
for i, (value, label) in enumerate(metrics):
    x = 7.0 + (i%2)*2.82; y = 1.48 + (i//2)*1.7
    rect(s, x, y, 2.55, 1.4, WHITE, True, LINE)
    textbox(s, value, x+.2, y+.2, 2.15, .42, 23, NAVY, True)
    textbox(s, label, x+.2, y+.75, 2.15, .38, 11, MUTED)
textbox(s, "DEFINITION OF DONE", .65, 4.15, 3.0, .24, 10, PINK, True)
textbox(s, "MVP selesai, proses bisnis tervalidasi, transaksi tercatat digital, dan pengguna pilot mampu mengakses layanan utama.", .65, 4.55, 11.8, .8, 20, NAVY, True)
footer(s, 4)

# 5 — Scope
s = base_slide("Project Scope", "04 · Boundaries")
card(s, "In Scope", ["Registrasi, login, profil & verifikasi WhatsApp", "Setor, tarik, saldo & histori simpanan", "Pembayaran biaya anggota", "Laporan individual & estimasi SHU", "Riwayat pembelian Lapak Employee"], .65, 1.45, 5.85, 4.65, GREEN, 18, 14)
card(s, "Out of Scope — Fase Awal", ["Pinjaman anggota", "Payroll deduction otomatis", "Integrasi penuh ERP rumah sakit", "Native mobile application", "Advanced analytics & loyalty engine"], 6.82, 1.45, 5.85, 4.65, RED, 18, 14)
footer(s, 5)

# 6 — Roadmap
s = base_slide("Delivery Roadmap", "05 · When & How")
steps = [("1", "Charter Approval", "H+0 · Non-MVP"), ("2", "Design", "H+21 · MVP ✓"), ("3", "Development", "H+77 · MVP ✓"), ("4", "UAT", "H+98 · MVP ✓"), ("5", "Pilot Go-Live", "H+120 · MVP ✓")]
rect(s, 1.05, 2.16, 11.15, .06, LINE, False, LINE)
for i, (n, title, desc) in enumerate(steps):
    x = .72 + i*2.48
    rect(s, x+.68, 1.78, .68, .68, PINK if i in (0,4) else NAVY2, True)
    textbox(s, n, x+.68, 1.98, .68, .22, 14, WHITE, True, PP_ALIGN.CENTER)
    textbox(s, title, x, 2.72, 2.05, .4, 14, NAVY, True, PP_ALIGN.CENTER)
    textbox(s, desc, x, 3.2, 2.05, .45, 10, MUTED, False, PP_ALIGN.CENTER)
for i, (title, desc) in enumerate([("Agile Delivery", "Sprint berprioritas nilai bisnis"), ("Prototype-First", "Validasi UX sebelum full build"), ("Pilot & Learn", "Ukur hasil lalu sempurnakan")]):
    card(s, title, desc, .75+i*4.2, 4.55, 3.8, 1.25, [PINK,PURPLE,NAVY2][i], 15, 11)
footer(s, 6)

# 7 — Governance and risks
s = base_slide("Governance & Key Risks", "06 · Control")
risks = [("Perubahan requirement", "PIC: PM & Product Owner", "TINGGI", RED), ("Adopsi pengguna rendah", "PIC: HR & Pengurus", "SEDANG", ORANGE), ("Keterlambatan integrasi", "PIC: IT & Tech Lead", "SEDANG", ORANGE), ("Keamanan data", "PIC: IT Security & Legal", "TINGGI", RED)]
for i, (risk, impact, level, color) in enumerate(risks):
    y = 1.45 + i*1.05
    rect(s, .65, y, 7.2, .82, WHITE, True, LINE)
    textbox(s, risk, .9, y+.18, 2.75, .25, 13, NAVY, True)
    textbox(s, impact, 3.75, y+.18, 2.05, .25, 11, MUTED)
    rect(s, 6.2, y+.16, 1.35, .38, color, True, color)
    textbox(s, level, 6.27, y+.27, 1.2, .14, 8, WHITE, True, PP_ALIGN.CENTER)
card(s, "Mitigasi Inti", ["Baseline scope & change control", "Onboarding dan feedback pengguna", "Mock service & integrasi bertahap", "Secure design, audit log, enkripsi"], 8.15, 1.45, 4.55, 4.0, PURPLE, 17, 13)
textbox(s, "dr. Mus Aida, MARS, MH → sumber daya & keputusan akhir   |   Mentor → bimbingan & arahan   |   Ketua/PM → memimpin & ikut mengerjakan   |   Tim → delivery harian", .75, 6.05, 11.9, .42, 10, NAVY, True, PP_ALIGN.CENTER)
footer(s, 7)

# 8 — Business value
s = base_slide("Business Value", "07 · Value Realization")
values = [("Transparansi", "Informasi dan transaksi anggota dapat ditelusuri"), ("Efisiensi", "Proses keanggotaan dan layanan lebih cepat"), ("Kesejahteraan", "Akses layanan koperasi yang relevan"), ("Pertumbuhan", "Fondasi unit usaha dan ekosistem digital")]
for i,(title,desc) in enumerate(values):
    x=.65+(i%2)*3.15; y=1.45+(i//2)*1.55
    card(s,title,desc,x,y,2.85,1.25,[PINK,PURPLE,NAVY2,GREEN][i],15,11)
card(s, "Investasi Awal · Rp 180–300 juta", "Design & development MVP, setup infrastruktur, security baseline, UAT, pelatihan, dan persiapan pilot.", 7.15, 1.45, 5.55, 1.25, RED, 14, 10)
card(s, "Operasi · Rp 12–25 juta/bulan", "Cloud, database & backup, WhatsApp/OTP, monitoring, maintenance, support, dan minor improvement.", 7.15, 2.95, 5.55, 1.35, ORANGE, 14, 10)
card(s, "Manfaat · Efisiensi 60–80%", "100% transaksi pilot digital; potensi GMV Lapak Employee Rp 30–75 juta/bulan. Estimasi indikatif, divalidasi melalui quotation dan pilot.", 7.15, 4.55, 5.55, 1.4, GREEN, 14, 10)
card(s, "Ritme Proyek (Cadence)", ["Check-in pertama: H+7", "Frekuensi: mingguan", "Pemimpin: May Dina — Ketua/PM", "Review: H+21, H+77, H+98, H+120"], .65, 4.7, 5.95, 1.35, PURPLE, 14, 9)
footer(s, 8)

# 9 — Stakeholder map
s = base_slide("Peta Stakeholder: Pengaruh × Dukungan", "08 · Engagement")
mx,my,mw,mh=.65,1.45,7.5,4.95
rect(s,mx,my,mw,mh,WHITE,True,LINE)
rect(s,mx+mw/2,my,.025,mh,LINE,False,LINE); rect(s,mx,my+mh/2,mw,.025,LINE,False,LINE)
textbox(s,"LAWAN BERPENGARUH",mx+.22,my+.18,2.6,.25,10,RED,True)
textbox(s,"CHAMPION",mx+mw-2.1,my+.18,1.8,.25,10,GREEN,True,PP_ALIGN.RIGHT)
textbox(s,"LAWAN LEMAH",mx+.22,my+mh-.48,2,.25,10,MUTED,True)
textbox(s,"FAN CLUB",mx+mw-1.8,my+mh-.48,1.5,.25,10,PURPLE,True,PP_ALIGN.RIGHT)
points=[(1,.77,.18),(2,.68,.30),(3,.61,.40),(4,.54,.23),(5,.70,.67),(6,.38,.33),(7,.29,.66),(8,.57,.76)]
for n,px,py in points:
    x=mx+px*mw; y=my+py*mh
    rect(s,x,y,.38,.38,PINK if n<5 else NAVY2,True)
    textbox(s,str(n),x,y+.1,.38,.12,8,WHITE,True,PP_ALIGN.CENTER)
names=["1  Direksi / Sponsor","2  HRD","3  Finance & Accounting","4  Information Technology","5  Pengurus Koperasi","6  Legal / Compliance","7  Perwakilan Karyawan","8  Tim KawanDigi"]
card(s,"Stakeholder",names,8.45,1.45,4.25,4.95,PURPLE,16,11)
footer(s, 9)

# 10 — Android overview
s = base_slide("Web Mobile KawanDigi", "LAMPIRAN · Product Experience", True)
textbox(s,"Satu aplikasi untuk layanan anggota",.65,1.3,5.2,.45,19,WHITE,True)
textbox(s,"Akses saldo, setoran, penarikan, kewajiban, transaksi, dan laporan keuangan dalam pengalaman mobile yang konsisten.",.65,1.9,4.5,1.1,14,RGBColor(220,210,235))
add_image_contain(s,ANDROID/"apps_main.png",5.65,1.25,2.3,5.75,WHITE)
add_image_contain(s,ANDROID/"app_login.png",8.15,1.55,1.8,4.75,WHITE)
add_image_contain(s,ANDROID/"apps_laporang_keuangan.png",10.15,1.55,2.25,4.75,WHITE)
for i,t in enumerate(["Secure login","Digital Saving Account","Laporan anggota"]):
    textbox(s,t,.65,3.45+i*.6,4.2,.3,13,PINK2 if i==1 else WHITE,True)
textbox(s,"10",12.15,7.12,.6,.2,9,RGBColor(190,180,210),True,PP_ALIGN.RIGHT)

# 11 — Android flows
s = base_slide("Key Member Flows", "LAMPIRAN · Web Mobile Screens")
screens=[("Setor", "apps_setor.png"),("Tarik", "apps_tarik.png"),("Pembayaran", "apps_pembayaran.png"),("Laporan", "apps_laporang_keuangan.png")]
for i,(label,file) in enumerate(screens):
    x=.58+i*3.17
    add_image_contain(s,ANDROID/file,x,1.35,2.75,4.95,WHITE)
    textbox(s,label,x,6.5,2.75,.3,12,NAVY,True,PP_ALIGN.CENTER)
footer(s, 11)

# 12 — Lapak employee overview
s = base_slide("LapakEmployee.com", "LAMPIRAN · Digital Business Unit")
textbox(s,"Marketplace internal sebagai unit usaha pertama KawanDigi",.65,1.35,6.0,.5,18,NAVY,True)
textbox(s,"Menghubungkan karyawan dengan seller internal melalui katalog produk, checkout, pembayaran, dan komunikasi WhatsApp.",.65,1.95,5.3,.9,13,MUTED)
add_image_contain(s,WEB/"lapakemployee.com.png",6.25,1.3,6.45,2.55,WHITE)
add_image_contain(s,WEB/"lapakemployee-view-product-seller.png",6.25,4.08,6.45,2.25,WHITE)
for i,t in enumerate(["Seller discovery","Product catalog","Checkout & payment","WhatsApp ordering"]):
    rect(s,.65,3.2+i*.72,4.9,.52,WHITE,True,LINE)
    textbox(s,f"{i+1:02d}",.82,3.36+i*.72,.4,.16,9,PINK,True)
    textbox(s,t,1.35,3.34+i*.72,3.9,.2,12,NAVY,True)
footer(s, 12)

# 13 — Seller and whatsapp
s = base_slide("Seller Operations & WhatsApp", "LAMPIRAN · Commerce Enablement")
add_image_contain(s,WEB/"lapakemployee-login-seller.png",.65,1.35,7.55,4.85,WHITE)
add_image_contain(s,WEB/"lapak-employee-whatsapp-order.png",8.45,1.35,4.25,4.85,WHITE)
textbox(s,"Seller Portal",.65,6.42,7.55,.28,12,NAVY,True,PP_ALIGN.CENTER)
textbox(s,"WhatsApp Order",8.45,6.42,4.25,.28,12,NAVY,True,PP_ALIGN.CENTER)
footer(s, 13)

# 14 — Commitment
s = prs.slides.add_slide(BLANK)
rect(s,0,0,13.333,7.5,NAVY,False)
rect(s,9.8,-.8,4.4,4.4,NAVY2,True,NAVY2); rect(s,-1.0,5.6,4.0,2.8,PINK,True,PINK)
textbox(s,"NEXT STEP",.75,.7,2.0,.25,11,PINK2,True)
textbox(s,"Alignment, Approval\n& Pilot Readiness",.75,1.2,8.0,1.25,31,WHITE,True)
textbox(s,"Komitmen bersama memastikan scope terkendali, kebutuhan bisnis tervalidasi, dan pilot memberikan evidence untuk keputusan berikutnya.",.75,2.85,7.25,.85,15,RGBColor(220,210,235))
for i,(role,action) in enumerate([("SPONSOR","Direktur Rumah Sakit"),("MENTOR","HOD Finance"),("KETUA PROYEK","Manajer Proyek"),("PERWAKILAN TIM","Anggota Inti")]):
    x=.55+i*3.18
    rect(s,x,4.45,2.9,1.2,RGBColor(69,35,116),True,RGBColor(102,68,150))
    textbox(s,role,x+.15,4.7,2.6,.25,12,WHITE,True,PP_ALIGN.CENTER)
    textbox(s,action,x+.15,5.13,2.6,.2,9,RGBColor(205,195,222),False,PP_ALIGN.CENTER)
textbox(s,"KawanDigi · Bersama Tumbuh, Bersama Sejahtera",.75,6.78,7.0,.22,10,WHITE,True)

prs.core_properties.title = "Project Charter — Inisiasi Koperasi KawanDigi - Brawijaya Hospital"
prs.core_properties.subject = "KawanDigi Project Charter and Product View"
prs.core_properties.author = "KawanDigi"
prs.core_properties.keywords = "KawanDigi, Project Charter, Koperasi Digital, Lapak Employee"
prs.save(OUT)
print(f"Created {OUT} ({len(prs.slides)} slides)")
